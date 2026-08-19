"""
infrastructure/local_storage/pptx_deck_store.py — PptxSession, the
python-pptx-backed in-memory/on-disk slide deck
AI Model Coder CLI v1.50.0 (Clean Architecture refactor, Phase C, Context #4)

Extracted 2026-08-18 from claude_powerpoint.py's PptxSession class.
Kept as one class rather than split method-by-method between domain/
and infrastructure/, per the CodeSession precedent from Context #3:
summary()/undo()/_snapshot()/apply_code()'s safety check are pure
logic, but apply_code() itself runs generated code against live
in-memory state (not meaningfully separable from the class it mutates)
and __init__/_load()/save()/_add_table()/_add_chart() all depend on
the python-pptx library and touch disk — the class as a whole belongs
here, not split apart.
"""

from domain.powerpoint import _DENYLIST

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None


class PptxSession:
    def __init__(self, input_path=None):
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for --pptx (pip install python-pptx)"
            )
        self.slides = []
        self._history_stack = []  # for /undo — list of deep-copied slide lists
        self._template_path = None  # if loaded from an existing deck, reuse its theme

        if input_path:
            self._load(input_path)

    def _load(self, path):
        prs = Presentation(path)
        self._template_path = path
        for slide in prs.slides:
            title = ""
            bullets = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if not title and shape == slide.shapes.title:
                    title = text
                elif text:
                    bullets.extend(line for line in text.split("\n") if line.strip())
            self.slides.append({
                "title": title, "bullets": bullets, "layout": "title_content",
                "table": None, "chart": None,
            })

    # ── context for the model ───────────────────────────────────────────

    def summary(self):
        if not self.slides:
            return "(no slides yet)"
        parts = [f"Deck has {len(self.slides)} slide(s):"]
        for i, s in enumerate(self.slides):
            extra = []
            if s.get("table"):
                extra.append(f"table {len(s['table']['rows'])}x{len(s['table']['headers'])}")
            if s.get("chart"):
                extra.append(f"{s['chart']['type']} chart")
            extra_str = f" [{', '.join(extra)}]" if extra else ""
            bullets_preview = "; ".join(s["bullets"][:3])
            parts.append(f"  {i}: \"{s['title']}\"{extra_str} — {bullets_preview}")
        return "\n".join(parts)

    # ── applying a model turn ───────────────────────────────────────────

    def _snapshot(self):
        import copy
        self._history_stack.append(copy.deepcopy(self.slides))
        if len(self._history_stack) > 20:
            self._history_stack.pop(0)

    def undo(self):
        if not self._history_stack:
            return False
        self.slides = self._history_stack.pop()
        return True

    def apply_code(self, code):
        """Run model-generated code against `self.slides`. Returns (ok, message)."""
        lowered = code.lower()
        for bad in _DENYLIST:
            if bad in lowered:
                return False, f"[blocked] generated code used a disallowed construct: {bad!r}"

        self._snapshot()

        def add_slide(title, bullets=None, layout="title_content", table=None, chart=None):
            self.slides.append({
                "title": title, "bullets": bullets or [], "layout": layout,
                "table": table, "chart": chart,
            })

        def update_slide(index, title=None, bullets=None, table=None, chart=None):
            s = self.slides[index]
            if title is not None:
                s["title"] = title
            if bullets is not None:
                s["bullets"] = bullets
            if table is not None:
                s["table"] = table
            if chart is not None:
                s["chart"] = chart

        def delete_slide(index):
            self.slides.pop(index)

        def reorder_slides(new_order):
            self.slides = [self.slides[i] for i in new_order]

        local_ns = {
            "slides": self.slides,
            "add_slide": add_slide,
            "update_slide": update_slide,
            "delete_slide": delete_slide,
            "reorder_slides": reorder_slides,
        }
        try:
            exec(compile(code, "<pptx-turn>", "exec"), {"__builtins__": {
                "len": len, "range": range, "sum": sum, "min": min, "max": max,
                "round": round, "sorted": sorted, "list": list, "dict": dict,
                "str": str, "int": int, "float": float, "bool": bool,
                "enumerate": enumerate, "zip": zip, "abs": abs,
            }}, local_ns)
        except Exception as e:
            self.undo()
            return False, f"[ERROR] generated code failed: {e}"
        return True, "applied"

    # ── persistence ──────────────────────────────────────────────────────

    def save(self, output_path):
        # Rebuilt from scratch on every save (simplest correct approach given
        # `slides` is the single source of truth) rather than diffing against
        # a live Presentation object — mirrors claude_excel.py rewriting the
        # whole workbook from `sheets` on every save.
        prs = Presentation()
        title_content = prs.slide_layouts[1]
        title_only = prs.slide_layouts[5]
        section_header = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else title_only

        for s in self.slides:
            layout = {"title_only": title_only, "section_header": section_header}.get(
                s.get("layout"), title_content)
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = s["title"]

            if s.get("bullets") and layout is title_content:
                body = None
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        body = ph
                        break
                if body is not None:
                    tf = body.text_frame
                    tf.clear()
                    for i, bullet in enumerate(s["bullets"]):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = bullet

            if s.get("table"):
                self._add_table(slide, s["table"])
            if s.get("chart"):
                self._add_chart(slide, s["chart"])

        prs.save(output_path)

    def _add_table(self, slide, table):
        headers, rows = table["headers"], table["rows"]
        n_rows, n_cols = len(rows) + 1, len(headers)
        left, top, width, height = Inches(0.5), Inches(1.8), Inches(9), Inches(0.4 * n_rows)
        shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        tbl = shape.table
        for c, header in enumerate(headers):
            tbl.cell(0, c).text = str(header)
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                tbl.cell(r, c).text = str(val)

    def _add_chart(self, slide, chart):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        xl_type = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }.get(chart["type"], XL_CHART_TYPE.COLUMN_CLUSTERED)

        data = CategoryChartData()
        data.categories = chart["categories"]
        for name, values in chart["series"].items():
            data.add_series(name, values)

        left, top, width, height = Inches(1), Inches(1.8), Inches(8), Inches(4.5)
        slide.shapes.add_chart(xl_type, left, top, width, height, data)
